"""Build the bilingual PPTX decks for the MAPAQ App Insights .NET 10 workshop.

Reads ``slides/content/<lang>/deck.yaml`` for both ``en`` and ``fr``, and writes:

  * ``docs/decks/app-insights-dotnet.pptx``     (en)
  * ``docs/decks/app-insights-dotnet-fr.pptx``  (fr)

All text is pinned to **Segoe UI** (DD-05 — Windows-safe). If a Mermaid PNG was
pre-rendered by ``slides/build/render-mermaid.mjs`` (filename pattern
``<slide-key>-<lang>.png``), the slide embeds it via ``slide.shapes.add_picture``
with ``descr=`` alt text for accessibility. If the template
``slides/theme/deck.pptx-template.pptx`` exists it is opened, otherwise the
script falls back to ``Presentation()``.
"""

from __future__ import annotations

import sys
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
CONTENT_DIR = ROOT / "slides" / "content"
TEMPLATE = ROOT / "slides" / "theme" / "deck.pptx-template.pptx"
ASSETS_DIR = ROOT / "docs" / "decks" / "assets"
OUT_DIR = ROOT / "docs" / "decks"

LANGS = ("en", "fr")
OUT_FILES = {
    "en": OUT_DIR / "app-insights-dotnet.pptx",
    "fr": OUT_DIR / "app-insights-dotnet-fr.pptx",
}

# Palette — matches slides/theme/deck.css
BG = RGBColor(0x0B, 0x12, 0x20)          # --bg
FG = RGBColor(0xE5, 0xE7, 0xEB)          # --fg
MUTED = RGBColor(0x94, 0xA3, 0xB8)       # muted
SKY = RGBColor(0x3B, 0x82, 0xF6)         # --sky
LEAF = RGBColor(0x16, 0xA3, 0x40)        # --leaf
ORANGE = RGBColor(0xF9, 0x73, 0x16)      # --orange
GRAPHITE = RGBColor(0x1F, 0x2A, 0x37)    # --graphite
CARD_BG = RGBColor(0x16, 0x1B, 0x22)     # --card-bg

FONT_NAME = "Segoe UI"  # DD-05 — Windows-safe

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def load_spec(lang: str) -> dict[str, Any]:
    path = CONTENT_DIR / lang / "deck.yaml"
    if not path.exists():
        raise FileNotFoundError(f"Missing deck spec: {path}")
    with path.open(encoding="utf-8") as f:
        return yaml.safe_load(f) or {}


def open_presentation() -> Presentation:
    if TEMPLATE.exists() and TEMPLATE.stat().st_size > 0:
        prs = Presentation(str(TEMPLATE))
    else:
        prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    return prs


def add_dark_background(slide, prs: Presentation) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    # Send to back so subsequent shapes render on top.
    sp_tree = bg._element.getparent()
    sp_tree.remove(bg._element)
    sp_tree.insert(2, bg._element)


def style_paragraph(p, *, size_pt: int, bold: bool = False, color: RGBColor = FG) -> None:
    p.font.name = FONT_NAME
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.color.rgb = color
    # Force every run inside the paragraph to the same font (defends against
    # python-pptx auto-creating runs that inherit from the master).
    for run in p.runs:
        run.font.name = FONT_NAME
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = color


def add_title(slide, text: str, *, size_pt: int = 32, color: RGBColor = FG) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    style_paragraph(tf.paragraphs[0], size_pt=size_pt, bold=True, color=color)


def add_subtitle(slide, text: str, *, size_pt: int = 20, top_in: float = 1.5) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top_in), Inches(12.1), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    style_paragraph(tf.paragraphs[0], size_pt=size_pt, color=MUTED)


def add_bullets(slide, bullets: list[str], *, top_in: float = 1.7, height_in: float = 4.5) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top_in), Inches(12.1), Inches(height_in))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        style_paragraph(p, size_pt=18, color=FG)


def add_picture_with_alt(slide, png_path: Path, alt_text: str, *, top_in: float = 2.2) -> None:
    pic = slide.shapes.add_picture(
        str(png_path), Inches(0.6), Inches(top_in), width=Inches(12.1)
    )
    if alt_text:
        pic._element.set("descr", alt_text)


def add_speaker_notes(slide, notes: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.text = notes
    style_paragraph(tf.paragraphs[0], size_pt=14, color=FG)


def find_diagram_png(slide_key: str, lang: str) -> Path | None:
    candidate = ASSETS_DIR / f"{slide_key}-{lang}.png"
    return candidate if candidate.exists() else None


def add_cover_slide(prs: Presentation, spec: dict[str, Any]) -> None:
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank)
    add_dark_background(slide, prs)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = spec.get("title", "")
    style_paragraph(tf.paragraphs[0], size_pt=44, bold=True, color=SKY)

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.0))
    sf = sub_box.text_frame
    sf.word_wrap = True
    sf.text = spec.get("subtitle", "")
    style_paragraph(sf.paragraphs[0], size_pt=22, color=MUTED)

    author = spec.get("author")
    if author:
        author_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.6))
        af = author_box.text_frame
        af.text = author
        style_paragraph(af.paragraphs[0], size_pt=14, color=LEAF)


def add_content_slide(prs: Presentation, slide_spec: dict[str, Any], lang: str) -> None:
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank)
    add_dark_background(slide, prs)

    add_title(slide, slide_spec.get("title", slide_spec.get("key", "")))

    bullets = slide_spec.get("bullets") or []
    has_diagram = bool(slide_spec.get("mermaid")) or bool(slide_spec.get("image"))

    if bullets and not has_diagram:
        add_bullets(slide, bullets, top_in=1.7, height_in=5.0)
    elif bullets and has_diagram:
        add_bullets(slide, bullets, top_in=1.5, height_in=1.5)

    slide_key = slide_spec.get("key", "")
    if slide_spec.get("mermaid"):
        png = find_diagram_png(slide_key, lang)
        alt = slide_spec.get("title") or slide_key
        if png is not None:
            top_in = 3.2 if bullets else 2.0
            add_picture_with_alt(slide, png, alt, top_in=top_in)
        else:
            # Diagram declared but PNG not pre-rendered — leave a placeholder
            # text so the missing asset is obvious in PowerPoint.
            box = slide.shapes.add_textbox(
                Inches(0.6), Inches(3.4), Inches(12.1), Inches(1.0)
            )
            tf = box.text_frame
            tf.text = f"[diagram pending: docs/decks/assets/{slide_key}-{lang}.png]"
            style_paragraph(tf.paragraphs[0], size_pt=14, color=ORANGE)

    if slide_spec.get("image"):
        img = ROOT / slide_spec["image"]
        if img.exists():
            top_in = 3.2 if bullets else 2.0
            add_picture_with_alt(slide, img, slide_spec.get("title", ""), top_in=top_in)

    notes = slide_spec.get("notes")
    if notes:
        add_speaker_notes(slide, notes)


def build_lang(lang: str) -> Path:
    spec = load_spec(lang)
    prs = open_presentation()

    add_cover_slide(prs, spec)

    for slide_spec in spec.get("slides", []):
        if slide_spec.get("key") == "cover":
            # Cover already rendered above; skip duplicate.
            continue
        add_content_slide(prs, slide_spec, lang)

    out = OUT_FILES[lang]
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def main() -> int:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    for lang in LANGS:
        try:
            out = build_lang(lang)
        except FileNotFoundError as exc:
            print(f"⚠ {exc}; skipping {lang}", file=sys.stderr)
            continue
        size_kb = out.stat().st_size / 1024
        print(f"✓ wrote {out} ({size_kb:.1f} KB)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
